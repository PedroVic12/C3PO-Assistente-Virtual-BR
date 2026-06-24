from abc import ABC, abstractmethod
from pptx import Presentation
from pptx.util import Inches
from markdown_pdf import MarkdownPdf
import os
import io
from PIL import Image, ImageDraw, ImageFont
import tempfile
import uuid
import sys
import pathlib

SRC_DIR = pathlib.Path(__file__).resolve().parent.parent
sys.path.append(str(SRC_DIR))


print(SRC_DIR)

from utils.parser import parser_singleton

# --- Interface (Strategy) ---
class Converter(ABC):
    @abstractmethod
    def convert(self, output_filename):
        pass

# --- Concrete Strategies ---
class PdfConverter(Converter):
    def convert(self, output_filename="apresentacao.pdf"):
        print("🔄 Convertendo para PDF...")
        md_slides = parser_singleton.get_slides_as_markdown()
        full_md_content = "\n\n<div style='page-break-after: always;'></div>\n\n".join(md_slides)
        
        pdf = MarkdownPdf()
        pdf.add_section(full_md_content)
        pdf.save(output_filename)
        print(f"✅ PDF gerado em: {output_filename}")

class PptxConverter(Converter):
    def convert(self, output_filename="apresentacao.pptx"):
        print("🔄 Convertendo para PowerPoint...")
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[5]
        
        md_slides = parser_singleton.get_slides_as_markdown()

        for md_slide in md_slides:
            title, content = parser_singleton.extract_title_and_content(md_slide)
            
            slide = prs.slides.add_slide(blank_slide_layout)
            
            if title:
                title_shape = slide.shapes.title
                title_shape.text = title

            if content:
                left = Inches(1)
                top = Inches(1.5)
                width = Inches(8)
                height = Inches(5)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.text = content
        
        prs.save(output_filename)
        print(f"✅ PowerPoint gerado em: {output_filename}")

class VideoConverter(Converter):
    def _create_image_from_text(self, text, size=(1280, 720)):
        """Cria uma imagem a partir de um texto de slide."""
        img = Image.new('RGB', size, color = (25, 25, 80))
        d = ImageDraw.Draw(img)
        
        try:
            font_title = ImageFont.truetype("arial.ttf", 60)
            font_content = ImageFont.truetype("arial.ttf", 30)
        except IOError:
            font_title = ImageFont.load_default()
            font_content = ImageFont.load_default()

        title, content = parser_singleton.extract_title_and_content(text)
        
        # Desenha o título
        d.text((size[0]*0.1, size[1]*0.1), title, fill=(255,255,255), font=font_title)
        
        # Desenha o conteúdo
        d.text((size[0]*0.1, size[1]*0.3), content, fill=(220,220,220), font=font_content)

        # Salva em um arquivo temporário
        temp_file = os.path.join(tempfile.gettempdir(), f"{uuid.uuid4()}.png")
        img.save(temp_file)
        return temp_file

    def convert(self, output_filename="apresentacao.mp4"):
        print("🔄 Convertendo para Vídeo (pode demorar)...")
        md_slides = parser_singleton.get_slides_as_markdown()
        clips = []
        temp_files = []
        try:

            for slide_md in md_slides:
                print(f"  -> Processando slide para vídeo...")
                img_path = self._create_image_from_text(slide_md)
                temp_files.append(img_path)
                clip = ImageClip(img_path).set_duration(2) # 2 segundos por slide
                clips.append(clip)
            
            if not clips:
                print("⚠️ Nenhum slide encontrado para gerar o vídeo.")
                return

            final_clip = concatenate_videoclips(clips, method="compose")
            final_clip.write_videofile(output_filename, fps=24, codec='libx264')

            # Limpa arquivos temporários
            for file_path in temp_files:
                os.remove(file_path)

            print(f"✅ Vídeo gerado em: {output_filename}")
            
        except Exception as e:
            print("Erro: ", e)
            for file_path in temp_files:
                os.remove(file_path)

# --- Factory Method ---
class ConverterFactory:
    def get_converter(self, format_type):
        if format_type.lower() == 'pdf':
            return PdfConverter()
        elif format_type.lower() == 'pptx':
            return PptxConverter()
        elif format_type.lower() == 'video':
            return VideoConverter()
        else:
            raise ValueError(f"Formato desconhecido: {format_type}")
